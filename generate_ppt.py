from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

# Colors
COLOR_1 = hex_to_rgb('#F76702') # Orange
COLOR_2 = hex_to_rgb('#F072AD') # Pink
COLOR_3 = hex_to_rgb('#642B87') # Purple
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(100, 100, 100)

prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Background shapes (Gradient-like effect using 3 bars)
left = top = Inches(0)
width = prs.slide_width
height = Inches(0.2)

# Decoration Bar at top
shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width/3, height)
shape1.fill.solid()
shape1.fill.fore_color.rgb = COLOR_1
shape1.line.fill.background()

shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + width/3, top, width/3, height)
shape2.fill.solid()
shape2.fill.fore_color.rgb = COLOR_2
shape2.line.fill.background()

shape3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + (width/3)*2, top, width/3, height)
shape3.fill.solid()
shape3.fill.fore_color.rgb = COLOR_3
shape3.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
p = title_tf.add_paragraph()
p.text = "Presentation Title"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = COLOR_3
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
subtitle_tf = subtitle_box.text_frame
p = subtitle_tf.add_paragraph()
p.text = "Subtitle or Presenter Name"
p.font.size = Pt(24)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# Bottom bar
bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.2), prs.slide_width, Inches(0.2))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = COLOR_3
bottom_bar.line.fill.background()


# Slide 2: Table of Contents / List
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Header Bar
header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), Inches(0.2), Inches(1))
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = COLOR_1
header_bar.line.fill.background()

# Header Text
header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
p = header_box.text_frame.add_paragraph()
p.text = "Table of Contents"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLACK

# Content List
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
tf = content_box.text_frame
for i, item in enumerate(["Introduction", "Key Concepts", "Analysis", "Conclusion"]):
    p = tf.add_paragraph()
    p.text = f"{i+1}. {item}"
    p.font.size = Pt(24)
    p.space_after = Pt(20)
    p.font.color.rgb = COLOR_3 if i % 2 == 0 else COLOR_2


# Slide 3: Content with Title
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Header
header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), Inches(0.2), Inches(1))
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = COLOR_2
header_bar.line.fill.background()

header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
p = header_box.text_frame.add_paragraph()
p.text = "Key Content Page"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLACK

# Body
body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8.5), Inches(4.5))
tf = body_box.text_frame
p = tf.add_paragraph()
p.text = "This is a clean content slide."
p.font.size = Pt(20)
p.space_after = Pt(14)

p = tf.add_paragraph()
p.text = "• Point 1: Highlighted in Orange"
p.font.size = Pt(18)
run = p.runs[0]
run.font.color.rgb = COLOR_1
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "• Point 2: Standard text description goes here. Keep it simple and clean."
p.font.size = Pt(18)
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "• Point 3: Highlighted in Pink"
p.font.size = Pt(18)
run = p.runs[0]
run.font.color.rgb = COLOR_2

# Slide 4: Section Header
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Full background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_3
bg.line.fill.background()

# Section Title
sec_title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
p = sec_title_box.text_frame.add_paragraph()
p.text = "SECTION 01"
p.font.size = Pt(20)
p.font.color.rgb = COLOR_1
p.alignment = PP_ALIGN.CENTER

p2 = sec_title_box.text_frame.add_paragraph()
p2.text = "Deep Dive"
p2.font.size = Pt(60)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER


prs.save('clean_template.pptx')
print("clean_template.pptx created successfully.")
